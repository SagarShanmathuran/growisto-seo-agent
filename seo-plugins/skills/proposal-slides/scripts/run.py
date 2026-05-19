"""
seo-proposal-slides – generate a 6-slide Ahrefs-powered proposal PPTX.

Two-pass workflow:

  Pass 1 (default):
    Load Ahrefs keyword + top-page CSVs → write data JSON for Claude to enrich.

  Pass 2 (--finalize --clusters-json PATH):
    Take Claude's enrichment JSON + base data → generate 6-slide PPTX.

Usage (Pass 1):
  python run.py \\
      --client-name Orra \\
      --client-csv ~/Downloads/orra-kw.csv \\
      --competitor "Tanishq:~/Downloads/tanishq-kw.csv" \\
      --competitor "Senco:~/Downloads/senco-kw.csv" \\
      [--top-pages-csv ~/Downloads/orra-pages.csv] \\
      [--comp-pages "Tanishq:~/Downloads/tanishq-pages.csv"] \\
      [--output-dir /tmp]

Usage (Pass 2):
  python run.py \\
      --client-name Orra --client-csv ~/Downloads/orra-kw.csv \\
      --competitor "Tanishq:..." --competitor "Senco:..." \\
      --clusters-json /tmp/clusters-orra.json \\
      --finalize [--output-dir ~/Downloads]
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path

import pandas as pd


# ─────────────────────────────────────────────────────────────────────────────
#  CSV loading  (handles UTF-16 tab and UTF-8 comma, Ahrefs formats)
# ─────────────────────────────────────────────────────────────────────────────

def _read_csv(path: str | Path) -> pd.DataFrame:
    path = Path(path).expanduser()
    for enc, sep in [("utf-16", "\t"), ("utf-8", ","), ("utf-8", "\t"), ("latin-1", ",")]:
        try:
            df = pd.read_csv(path, encoding=enc, sep=sep)
            if df.shape[1] > 1:
                return df
        except Exception:
            continue
    raise ValueError(f"Cannot parse CSV: {path}")


_ALIASES: dict[str, str] = {
    "keyword": "Keyword",
    "current position": "Position",
    "position": "Position",
    "volume": "Volume",
    "search volume": "Volume",
    "top keyword: volume": "Volume",
    "top keyword volume": "Volume",
    "current organic traffic": "Traffic",
    "organic traffic": "Traffic",
    "traffic": "Traffic",
    "current url": "URL",
    "url": "URL",
    "page url": "URL",
    "top keyword": "Keyword",
    "top keyword: position": "Position",
}


def _normalise(df: pd.DataFrame) -> pd.DataFrame:
    rename: dict[str, str] = {}
    for col in df.columns:
        canon = _ALIASES.get(col.lower().strip())
        if canon and canon not in df.columns:
            rename[col] = canon
    df = df.rename(columns=rename)
    for c in ("Position", "Volume", "Traffic"):
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    if "Keyword" in df.columns:
        df = df[df["Keyword"].astype(str).str.strip().str.len() > 0].copy()
    return df


def load_kw(path: str | Path) -> pd.DataFrame:
    return _normalise(_read_csv(path))


# ─────────────────────────────────────────────────────────────────────────────
#  Formatting helpers
# ─────────────────────────────────────────────────────────────────────────────

BUCKETS = ["1", "2", "3", "(4-6)", "(7-10)", "(11-20)", "(21-30)", "(31-50)", "Not Ranking"]


def _bucket(pos) -> str:
    if pd.isna(pos) or pos <= 0:
        return "Not Ranking"
    pos = int(pos)
    if pos == 1: return "1"
    if pos == 2: return "2"
    if pos == 3: return "3"
    if pos <= 6:  return "(4-6)"
    if pos <= 10: return "(7-10)"
    if pos <= 20: return "(11-20)"
    if pos <= 30: return "(21-30)"
    if pos <= 50: return "(31-50)"
    return "Not Ranking"


def _fmt_pos(pos) -> str:
    if pd.isna(pos) or pos <= 0 or pos > 100:
        return "NR"
    return str(int(pos))


def _fmt_num(n: int | float) -> str:
    n = int(n)
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.1f}K"
    return str(n)


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 1 – Keyword ranking comparison
# ─────────────────────────────────────────────────────────────────────────────

def compute_slide1(
    client_df: pd.DataFrame,
    comp_dfs: list[pd.DataFrame],
    comp_names: list[str],
    n: int = 15,
) -> dict:
    if "Keyword" not in client_df.columns:
        return {"rows": [], "comp_names": comp_names}

    client_map: dict[str, dict] = {}
    for _, row in client_df.iterrows():
        kw = str(row["Keyword"]).strip().lower()
        client_map[kw] = {
            "pos": int(row["Position"]) if row["Position"] > 0 else 999,
            "vol": int(row.get("Volume", 0)),
        }

    comp_maps: list[dict[str, int]] = []
    for df in comp_dfs:
        m: dict[str, int] = {}
        if "Keyword" in df.columns:
            for _, row in df.iterrows():
                kw = str(row["Keyword"]).strip().lower()
                pos = int(row["Position"]) if row["Position"] > 0 else 999
                m[kw] = pos
        comp_maps.append(m)

    # Candidates: any competitor ranks top-10 AND client ranks worse
    candidates: dict[str, int] = {}  # kw -> best available volume
    for comp_map, df in zip(comp_maps, comp_dfs):
        for kw, comp_pos in comp_map.items():
            if comp_pos > 10:
                continue
            client_pos = client_map.get(kw, {}).get("pos", 999)
            if client_pos <= comp_pos:
                continue
            # pick best volume from client or competitor
            vol = client_map.get(kw, {}).get("vol", 0)
            if "Keyword" in df.columns and "Volume" in df.columns:
                rows = df[df["Keyword"].str.lower().str.strip() == kw]
                if not rows.empty:
                    cv = int(rows.iloc[0]["Volume"])
                    if cv > vol:
                        vol = cv
            if vol > candidates.get(kw, -1):
                candidates[kw] = vol

    rows = []
    for kw in sorted(candidates, key=lambda k: candidates[k], reverse=True)[:n]:
        client_pos_str = _fmt_pos(client_map.get(kw, {}).get("pos", 0))
        r: dict = {
            "keyword": kw.title(),
            "volume": _fmt_num(candidates[kw]),
            "client_pos": client_pos_str,
        }
        for i, comp_map in enumerate(comp_maps):
            r[f"comp{i}_pos"] = _fmt_pos(comp_map.get(kw, 0))
        rows.append(r)

    return {"rows": rows, "comp_names": comp_names}


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 2 – Ranking bucket distribution
# ─────────────────────────────────────────────────────────────────────────────

def compute_slide2(client_df: pd.DataFrame, client_name: str) -> dict:
    if "Keyword" not in client_df.columns:
        return {"buckets": [], "client_name": client_name}

    df = client_df.copy()
    df["_bucket"] = df["Position"].apply(_bucket)

    grouped = df.groupby("_bucket").agg(
        kw_count=("Keyword", "count"),
        total_sv=("Volume", "sum"),
    ).reset_index()

    total_sv = df["Volume"].sum()

    result_buckets = []
    for b in BUCKETS:
        r = grouped[grouped["_bucket"] == b]
        if r.empty:
            result_buckets.append({"bucket": b, "kw_count": 0, "total_sv": 0, "pct_sv": 0.0})
        else:
            sv = int(r.iloc[0]["total_sv"])
            result_buckets.append({
                "bucket": b,
                "kw_count": int(r.iloc[0]["kw_count"]),
                "total_sv": sv,
                "pct_sv": round(sv / total_sv * 100, 2) if total_sv else 0.0,
            })

    result_buckets.append({
        "bucket": "Grand Total",
        "kw_count": int(df["Keyword"].count()),
        "total_sv": int(total_sv),
        "pct_sv": 100.0,
    })

    return {"buckets": result_buckets, "client_name": client_name}


def _slide2_insights(buckets: list[dict]) -> list[str]:
    """Auto-generate insight bullets from bucket data."""
    mid_sv = sum(b["total_sv"] for b in buckets if b["bucket"] in ("(4-6)", "(7-10)", "(11-20)"))
    total_sv = next((b["total_sv"] for b in buckets if b["bucket"] == "Grand Total"), 1) or 1
    mid_pct = round(mid_sv / total_sv * 100, 1)

    def _pct(bucket_name: str) -> float:
        return next((b["pct_sv"] for b in buckets if b["bucket"] == bucket_name), 0.0)

    insights = []
    if mid_pct:
        insights.append(
            f"Mid-rankings (4-20) dominate opportunity: ~{mid_pct}% of total search volume sits here, "
            "indicating strong potential to drive growth by improving existing rankings."
        )
    p7_10 = _pct("(7-10)")
    if p7_10:
        insights.append(
            f"Positions 7-10 ({p7_10}%) highlight a large set of near-page-1 keywords "
            "that can be pushed higher with targeted optimization."
        )
    p4_6 = _pct("(4-6)")
    if p4_6:
        insights.append(
            f"Positions 4-6 ({p4_6}%) act as quick wins — movement into top 3 unlocks "
            "significantly higher click-through rates."
        )
    p11_20 = _pct("(11-20)")
    if p11_20:
        insights.append(
            f"Positions 11-20 ({p11_20}%) offer scalable gains through content depth "
            "and authority improvements."
        )
    return insights[:4]


# ─────────────────────────────────────────────────────────────────────────────
#  Data for Claude (slides 3-6)
# ─────────────────────────────────────────────────────────────────────────────

def summarise_for_claude(
    client_df: pd.DataFrame,
    comp_dfs: list[pd.DataFrame],
    client_name: str,
    comp_names: list[str],
    top_pages_map: dict[str, pd.DataFrame],
) -> dict:
    """Extract keyword + page summaries for Claude to cluster and categorize."""

    client_kws = []
    if "Keyword" in client_df.columns:
        for _, row in client_df.sort_values("Volume", ascending=False).head(300).iterrows():
            kw = str(row["Keyword"]).strip()
            if kw:
                client_kws.append({
                    "kw": kw,
                    "vol": int(row.get("Volume", 0)),
                    "pos": _fmt_pos(row.get("Position", 0)),
                    "traffic": int(row.get("Traffic", 0)),
                    "url": str(row.get("URL", "")),
                })

    comp_kws: dict[str, list] = {}
    for name, df in zip(comp_names, comp_dfs):
        if "Keyword" not in df.columns:
            continue
        comp_kws[name] = []
        for _, row in df.sort_values("Volume", ascending=False).head(200).iterrows():
            kw = str(row["Keyword"]).strip()
            if kw:
                comp_kws[name].append({
                    "kw": kw,
                    "vol": int(row.get("Volume", 0)),
                    "pos": _fmt_pos(row.get("Position", 0)),
                    "traffic": int(row.get("Traffic", 0)),
                    "url": str(row.get("URL", "")),
                })

    top_pages: dict[str, list] = {}
    for domain, df in top_pages_map.items():
        if "URL" not in df.columns:
            continue
        sort_col = "Traffic" if "Traffic" in df.columns else df.columns[1]
        top_pages[domain] = []
        for _, row in df.sort_values(sort_col, ascending=False).head(40).iterrows():
            entry: dict = {"url": str(row.get("URL", "")), "traffic": int(row.get("Traffic", 0))}
            if "Keyword" in df.columns:
                entry["top_kw"] = str(row["Keyword"])
            if "Volume" in df.columns:
                entry["kw_sv"] = int(row["Volume"])
            top_pages[domain].append(entry)

    return {
        "client_name": client_name,
        "comp_names": comp_names,
        "client_top_keywords": client_kws,
        "competitor_keywords": comp_kws,
        "top_pages": top_pages,
    }


# ─────────────────────────────────────────────────────────────────────────────
#  PPTX helpers  —  10 × 5.625"  (matches reference deck format)
# ─────────────────────────────────────────────────────────────────────────────

_SW = 10.0    # slide width  (inches)
_SH = 5.625   # slide height (inches)
_IN = 914400  # EMU per inch


def _i(x: float) -> int:
    """Inches → EMU."""
    return int(x * _IN)


def _make_prs():
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width  = _i(_SW)
    prs.slide_height = _i(_SH)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# ── Palette (matches reference deck) ─────────────────────────────────────────
TEAL        = lambda: _rgb(0x36, 0x75, 0x88)
DARK_NAVY   = lambda: _rgb(0x1C, 0x1C, 0x2E)
WHITE       = lambda: _rgb(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = lambda: _rgb(0x21, 0x21, 0x21)
MID_GRAY    = lambda: _rgb(0x4A, 0x55, 0x68)
LIGHT_GRAY  = lambda: _rgb(0xFA, 0xFC, 0xFD)
ALT_ROW     = lambda: _rgb(0xF5, 0xF8, 0xFA)
BORDER_CLR  = lambda: _rgb(0xDD, 0xE3, 0xE8)
ORANGE      = lambda: _rgb(0xE8, 0x87, 0x2A)
RED_ALERT   = lambda: _rgb(0xC0, 0x39, 0x2B)
GREEN_OK    = lambda: _rgb(0x28, 0xA7, 0x45)
TEAL_TINT   = lambda: _rgb(0xE8, 0xF3, 0xF5)
ORANGE_TINT = lambda: _rgb(0xFE, 0xF3, 0xE8)
RED_TINT    = lambda: _rgb(0xFD, 0xEC, 0xEA)
WARM_PINK   = lambda: _rgb(0xFF, 0xF8, 0xF5)


def _pos_colors(pos_str: str):
    if pos_str in ("NR", "", "Not Ranking"):
        return _rgb(0xC0, 0x39, 0x2B), WHITE()
    try:
        p = int(pos_str)
        if p <= 3:  return _rgb(0x1A, 0x73, 0x48), WHITE()
        if p <= 10: return _rgb(0xC7, 0x72, 0x00), WHITE()
        if p <= 20: return _rgb(0xF9, 0xA8, 0x25), NEAR_BLACK()
        if p <= 50: return _rgb(0xE6, 0x51, 0x00), WHITE()
        return _rgb(0xC0, 0x39, 0x2B), WHITE()
    except (ValueError, TypeError):
        return _rgb(0xC0, 0x39, 0x2B), WHITE()


def _gap_colors(gap_str: str):
    """Return (bg_rgb, text_rgb) for a gap multiple string like '17x' or '>100x'."""
    nums = re.findall(r'\d+', gap_str)
    try:
        n = float(nums[0])
        if n >= 50 or ">" in gap_str:
            return RED_ALERT(), WHITE()
        if n >= 20:
            return _rgb(0xE6, 0x51, 0x00), WHITE()
        return ORANGE(), WHITE()
    except (IndexError, ValueError):
        return ORANGE(), WHITE()


def _rect(slide, left, top, width, height, fill, border=None):
    s = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = _i(0.007)
    else:
        s.line.fill.background()
    return s


def _hline(slide, left, top, width, color=None):
    _rect(slide, int(left), int(top), int(width), _i(0.01), color or BORDER_CLR())


def _textbox(slide, left, top, width, height, text, size=9, bold=False,
             color=None, align="left", wrap=True):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    if color is None:
        color = NEAR_BLACK()
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _textbox_lines(slide, left, top, width, height, lines: list,
                   size=8.5, bold=False, color=None, wrap=True):
    """Multi-paragraph textbox — avoids the \\n-in-run rendering bug."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if color is None:
        color = NEAR_BLACK()
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = str(line)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def _cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def _cell_text(cell, text: str, size: float = 8.5, bold: bool = False,
               color=None, align: str = "center"):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _ALIGN = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}
    if color is None:
        color = NEAR_BLACK()
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(align, PP_ALIGN.CENTER)
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    cell.margin_left   = Pt(3)
    cell.margin_right  = Pt(3)
    cell.margin_top    = Pt(2)
    cell.margin_bottom = Pt(2)


def _make_table(slide, data_rows, headers, col_widths, left, top,
                row_h, header_h=None):
    if header_h is None:
        header_h = int(row_h * 1.6)
    nrows = len(data_rows) + 1
    ncols = len(headers)
    total_w = int(sum(col_widths))
    total_h = int(header_h + row_h * len(data_rows))
    tbl = slide.shapes.add_table(nrows, ncols, int(left), int(top),
                                  total_w, total_h).table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = int(w)
    tbl.rows[0].height = int(header_h)
    for i in range(1, nrows):
        tbl.rows[i].height = int(row_h)
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        _cell_fill(cell, TEAL())
        _cell_text(cell, h, size=8, bold=True, color=WHITE(), align="center")
    for i, row_vals in enumerate(data_rows):
        bg = ALT_ROW() if i % 2 == 0 else WHITE()
        for j, val in enumerate(row_vals):
            cell = tbl.cell(i + 1, j)
            if isinstance(val, tuple):
                _cell_fill(cell, val[1])
                _cell_text(cell, val[0], size=8.5, color=val[2], align="center")
            else:
                _cell_fill(cell, bg)
                al = "left" if j == 0 else "center"
                _cell_text(cell, str(val), size=8.5, color=NEAR_BLACK(), align=al)
    return tbl


def _header(slide, title: str, subtitle: str = ""):
    W = _i(_SW)
    H = _i(0.58) if subtitle else _i(0.48)
    _rect(slide, 0, 0, W, H, TEAL())
    if subtitle:
        _textbox(slide, _i(0.28), _i(0.05), W - _i(0.56), _i(0.28),
                 title, size=11, bold=True, color=WHITE())
        _textbox(slide, _i(0.28), _i(0.33), W - _i(0.56), _i(0.22),
                 subtitle, size=7.5, color=_rgb(0xCA, 0xE8, 0xF0))
    else:
        _textbox(slide, _i(0.28), _i(0.1), W - _i(0.56), _i(0.35),
                 title, size=11, bold=True, color=WHITE())
    _hline(slide, 0, H, W)


def _footer(slide, text: str):
    _textbox(slide, _i(0.28), _i(5.44), _i(9.44), _i(0.18),
             text, size=7, color=MID_GRAY())


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 1  –  Keyword Ranking: Client vs Competitors
# ─────────────────────────────────────────────────────────────────────────────

def slide1_keyword_ranking(prs, data: dict, date_str: str, client_name: str):
    slide = _blank(prs)
    comp_names: list = data.get("comp_names", [])
    rows       = data.get("rows", [])
    n_comp     = min(len(comp_names), 2)

    comp_str = " & ".join(comp_names[:2]) if comp_names else "competitors"
    _header(slide,
            f"Significant potential for {client_name} to improve rankings vs {comp_str}",
            f"Keywords where {comp_str} outrank {client_name} — sorted by monthly search volume  |  Source: Ahrefs, {date_str}")
    _footer(slide, f"Source: Ahrefs as of {date_str}.  NR = Not Ranking.  Snapshot; full keyword list finalised post content strategy.")

    if not rows:
        _textbox(slide, _i(0.3), _i(2.0), _i(9.4), _i(0.5),
                 "No keyword data found. Check CSV paths.", size=10, color=MID_GRAY())
        return

    AVAIL_W  = _i(9.44)
    KW_W     = _i(3.2)
    SV_W     = _i(1.1)
    pos_w    = (AVAIL_W - KW_W - SV_W) / (1 + n_comp)
    col_widths = [KW_W, SV_W] + [int(pos_w)] * (1 + n_comp)
    headers  = ["Keyword", "Avg. Monthly SV", client_name] + comp_names[:n_comp]

    table_rows = []
    for r in rows:
        bg_c, fg_c = _pos_colors(r["client_pos"])
        row_data = [r["keyword"], r["volume"], (r["client_pos"], bg_c, fg_c)]
        for i in range(n_comp):
            p = r.get(f"comp{i}_pos", "NR")
            bg, fg = _pos_colors(p)
            row_data.append((p, bg, fg))
        table_rows.append(row_data)

    TABLE_TOP = _i(0.62)
    _make_table(slide, table_rows, headers, col_widths,
                left=_i(0.28), top=TABLE_TOP, row_h=_i(0.245), header_h=_i(0.40))

    table_h = _i(0.40) + len(table_rows) * _i(0.245)
    callout_top = TABLE_TOP + table_h + _i(0.1)

    # Callout insight box
    callout_h = _i(0.36)
    _rect(slide, _i(0.28), callout_top, _i(9.44), callout_h, WHITE(), border=BORDER_CLR())
    msg = (f"Several keywords above show {client_name} ranking well below competitors — "
           "closing these gaps through on-page optimisation and content depth is the primary opportunity.")
    _textbox(slide, _i(0.38), callout_top + _i(0.05), _i(9.24), callout_h - _i(0.08),
             msg, size=8, color=MID_GRAY())

    # Legend
    legend_top = callout_top + callout_h + _i(0.06)
    _textbox(slide, _i(0.28), legend_top, _i(9.44), _i(0.2),
             "Position:  Green = 1-3  |  Amber = 4-10  |  Yellow = 11-20  |  Orange = 21-50  |  Red = NR",
             size=7, color=MID_GRAY())


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 2  –  Ranking Bucketing
# ─────────────────────────────────────────────────────────────────────────────

def slide2_ranking_bucket(prs, data: dict, date_str: str):
    slide = _blank(prs)
    client_name = data.get("client_name", "Client")
    buckets     = data.get("buckets", [])

    _header(slide,
            f"Unlocking Growth Opportunities Across Mid & Unranked Keyword Buckets — {client_name}",
            "Position distribution by keyword count and search volume share")
    _footer(slide, f"Source: Ahrefs as of {date_str}.")

    mid_sv   = sum(b["total_sv"] for b in buckets if b["bucket"] in ("(4-6)", "(7-10)", "(11-20)"))
    total_sv = next((b["total_sv"] for b in buckets if b["bucket"] == "Grand Total"), 1) or 1
    mid_pct  = round(mid_sv / total_sv * 100, 1)
    nr_pct   = next((b["pct_sv"] for b in buckets if b["bucket"] == "Not Ranking"), 0.0)

    # ── Left stat panel ───────────────────────────────────────────────────────
    P_LEFT = _i(0.22)
    P_TOP  = _i(0.62)
    P_W    = _i(2.1)
    P_H    = _i(4.5)
    _rect(slide, P_LEFT, P_TOP, P_W, P_H, TEAL_TINT(), border=BORDER_CLR())

    _textbox(slide, P_LEFT + _i(0.08), P_TOP + _i(0.18), P_W - _i(0.16), _i(0.95),
             f"{mid_pct}%", size=40, bold=True, color=TEAL(), align="center")
    _textbox(slide, P_LEFT + _i(0.1), P_TOP + _i(1.15), P_W - _i(0.2), _i(0.75),
             "of search volume\nin mid-rankings\n(positions 4–20)",
             size=8, color=DARK_NAVY(), align="center", wrap=True)

    _hline(slide, P_LEFT + _i(0.2), P_TOP + _i(2.1), P_W - _i(0.4))

    _textbox(slide, P_LEFT + _i(0.08), P_TOP + _i(2.22), P_W - _i(0.16), _i(0.55),
             f"{nr_pct:.1f}%", size=24, bold=True, color=RED_ALERT(), align="center")
    _textbox(slide, P_LEFT + _i(0.1), P_TOP + _i(2.82), P_W - _i(0.2), _i(0.5),
             "not yet ranking —\nuntapped demand",
             size=7.5, color=MID_GRAY(), align="center", wrap=True)

    # ── Right table ───────────────────────────────────────────────────────────
    TBL_LEFT = P_LEFT + P_W + _i(0.18)
    TBL_TOP  = P_TOP
    TBL_W    = _i(_SW) - TBL_LEFT - _i(0.22)

    raw_widths  = [_i(1.5), _i(1.1), _i(1.5), _i(1.0)]
    scale       = TBL_W / sum(raw_widths)
    col_widths  = [int(w * scale) for w in raw_widths]

    headers = ["Position Range", "Keywords", "Search Volume", "% of SV"]
    table_rows = []
    for b in buckets:
        is_total = b["bucket"] == "Grand Total"
        sv_str   = _fmt_num(b["total_sv"])
        pct_str  = f"{b['pct_sv']:.1f}%"
        if is_total:
            row = [
                ("Grand Total", DARK_NAVY(), WHITE()),
                (str(b["kw_count"]), DARK_NAVY(), WHITE()),
                (sv_str, DARK_NAVY(), WHITE()),
                (pct_str, DARK_NAVY(), WHITE()),
            ]
        else:
            row = [b["bucket"], str(b["kw_count"]), sv_str, pct_str]
        table_rows.append(row)

    _make_table(slide, table_rows, headers, col_widths,
                left=TBL_LEFT, top=TBL_TOP, row_h=_i(0.37), header_h=_i(0.42))

    # Red/orange left-accent highlights on rows (4-6) and (7-10)
    BUCKET_ORDER = ["1", "2", "3", "(4-6)", "(7-10)", "(11-20)", "(21-30)", "(31-50)",
                    "Not Ranking", "Grand Total"]
    HEADER_H_EMU = _i(0.42)
    ROW_H_EMU    = _i(0.37)
    for bucket_name, accent_color in [("(4-6)", ORANGE()), ("(7-10)", RED_ALERT())]:
        try:
            data_idx  = BUCKET_ORDER.index(bucket_name)  # 0-based position in data
            row_top   = TBL_TOP + HEADER_H_EMU + data_idx * ROW_H_EMU
            _rect(slide, TBL_LEFT, row_top, _i(0.07), ROW_H_EMU, accent_color)
        except ValueError:
            pass

    # Insight bullets below the panel/table
    tbl_bottom = TBL_TOP + _i(0.42) + len(table_rows) * _i(0.37)
    bullet_top = max(P_TOP + P_H, tbl_bottom) + _i(0.1)
    insights   = _slide2_insights(buckets)
    if bullet_top < _i(5.2) and insights:
        lines = [f"•  {b}" for b in insights[:3]]
        _textbox_lines(slide, _i(0.22), bullet_top, _i(9.56), _i(0.85),
                       lines, size=8, color=MID_GRAY())


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 3  –  Topical Clusters (hub-and-spoke ≤4; grid 5-6)
# ─────────────────────────────────────────────────────────────────────────────

def slide3_topical_cluster(prs, clusters_data: dict, date_str: str, client_name: str):
    slide = _blank(prs)
    clusters: list = clusters_data.get("slide3_clusters", [])

    _header(slide,
            f"Topical Clusters — The Content Universe We're Building for {client_name}",
            f"GKP-verified keyword demand grouped into interlinked content clusters  |  {date_str}")
    _footer(slide, f"Data: Ahrefs Organic Keywords + LLM clustering  |  {date_str}")

    if not clusters:
        _textbox(slide, _i(0.3), _i(2.5), _i(9.4), _i(0.5),
                 "No cluster data. Populate slide3_clusters in clusters JSON.", size=10, color=MID_GRAY())
        return

    n         = min(len(clusters), 6)
    HEADER_H  = _i(0.60)
    FOOTER_Y  = _i(5.44)
    AVAIL_H   = FOOTER_Y - HEADER_H

    if n <= 4:
        # ── Hub-and-spoke layout ──────────────────────────────────────────────
        HUB_W  = _i(2.9)
        HUB_H  = _i(0.72)
        HUB_L  = _i((_SW - 2.9) / 2)
        HUB_T  = HEADER_H + (AVAIL_H - HUB_H) // 2
        HUB_CX = HUB_L + HUB_W // 2
        HUB_CY = HUB_T + HUB_H // 2

        _rect(slide, HUB_L, HUB_T, HUB_W, HUB_H, DARK_NAVY(), border=TEAL())
        _textbox(slide, HUB_L + _i(0.1), HUB_T + _i(0.1),
                 HUB_W - _i(0.2), HUB_H - _i(0.18),
                 f"{client_name}  Content Hub",
                 size=9.5, bold=True, color=WHITE(), align="center")

        CARD_W = _i(2.9)
        CARD_H = _i(1.92)
        card_positions = [
            (_i(0.18), HEADER_H + _i(0.05)),                              # TL
            (_i(_SW) - CARD_W - _i(0.18), HEADER_H + _i(0.05)),           # TR
            (_i(0.18), FOOTER_Y - CARD_H - _i(0.06)),                     # BL
            (_i(_SW) - CARD_W - _i(0.18), FOOTER_Y - CARD_H - _i(0.06)), # BR
        ]
        card_anchors = [
            (card_positions[0][0] + CARD_W, card_positions[0][1] + CARD_H // 2),
            (card_positions[1][0],           card_positions[1][1] + CARD_H // 2),
            (card_positions[2][0] + CARD_W, card_positions[2][1] + CARD_H // 2),
            (card_positions[3][0],           card_positions[3][1] + CARD_H // 2),
        ]
        hub_anchors = [
            (HUB_L,          HUB_CY),
            (HUB_L + HUB_W, HUB_CY),
            (HUB_L,          HUB_CY),
            (HUB_L + HUB_W, HUB_CY),
        ]
        try:
            from pptx.enum.shapes import MSO_CONNECTOR_TYPE
            from pptx.util import Pt
            for (hx, hy), (ax, ay) in zip(hub_anchors, card_anchors):
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT, hx, hy, ax, ay)
                conn.line.color.rgb = ORANGE()
                conn.line.width = Pt(1.5)
        except Exception:
            pass

        for idx, cl in enumerate(clusters[:n]):
            cl_l, cl_t = card_positions[idx]
            _draw_cluster_card(slide, cl_l, cl_t, CARD_W, CARD_H, cl, idx + 1)

    else:
        # ── Hub title bar + grid ──────────────────────────────────────────────
        HUB_H = _i(0.42)
        _rect(slide, _i(0.22), HEADER_H + _i(0.06), _i(9.56), HUB_H, DARK_NAVY())
        _textbox(slide, _i(0.22), HEADER_H + _i(0.06), _i(9.56), HUB_H,
                 f"{client_name} Content Hub  •  {n} Focus Clusters",
                 size=10, bold=True, color=WHITE(), align="center")

        COLS      = 3
        ROWS_CNT  = (n + COLS - 1) // COLS
        GAP_X     = _i(0.12)
        GAP_Y     = _i(0.1)
        GRID_TOP  = HEADER_H + HUB_H + _i(0.1)
        GRID_W    = _i(9.56)
        GRID_H    = FOOTER_Y - GRID_TOP - _i(0.06)
        CARD_W    = (GRID_W - (COLS - 1) * GAP_X) // COLS
        CARD_H    = (GRID_H - (ROWS_CNT - 1) * GAP_Y) // ROWS_CNT

        for idx, cl in enumerate(clusters[:n]):
            col_i  = idx % COLS
            row_i  = idx // COLS
            card_l = _i(0.22) + col_i * (CARD_W + GAP_X)
            card_t = GRID_TOP + row_i * (CARD_H + GAP_Y)
            _draw_cluster_card(slide, card_l, card_t, CARD_W, CARD_H, cl, idx + 1)


def _draw_cluster_card(slide, left, top, width, height, cl: dict, num: int):
    """Render a single cluster card: border, number badge, metrics chip, keyword chips."""
    name    = cl.get("name", "Cluster")
    kw_cnt  = cl.get("kw_count", 0)
    tot_sv  = cl.get("total_sv", 0)
    samples = cl.get("sample_keywords", [])[:3]
    desc    = cl.get("description", "")

    PAD = _i(0.11)
    _rect(slide, left, top, width, height, LIGHT_GRAY(), border=BORDER_CLR())

    # Number badge
    BADGE_SZ = _i(0.26)
    _rect(slide, left + PAD, top + PAD, BADGE_SZ, BADGE_SZ, TEAL())
    _textbox(slide, left + PAD, top + PAD, BADGE_SZ, BADGE_SZ,
             f"{num:02d}", size=7.5, bold=True, color=WHITE(), align="center")

    # Cluster name
    _textbox(slide, left + PAD + BADGE_SZ + _i(0.07), top + PAD,
             width - PAD * 2 - BADGE_SZ - _i(0.07), BADGE_SZ,
             name, size=9, bold=True, color=DARK_NAVY())

    y = top + PAD + BADGE_SZ + _i(0.1)

    # Metrics row (teal-tint chip)
    M_H = _i(0.27)
    _rect(slide, left + PAD, y, width - PAD * 2, M_H, TEAL_TINT())
    _textbox(slide, left + PAD + _i(0.05), y + _i(0.04),
             width - PAD * 2 - _i(0.1), M_H - _i(0.06),
             f"{_fmt_num(tot_sv)} total SV   •   {kw_cnt} keywords",
             size=7.5, bold=True, color=TEAL())
    y += M_H + _i(0.1)

    # Keyword chips
    if samples:
        chip_x = left + PAD
        chip_y = y
        CHIP_H = _i(0.21)
        for kw in samples:
            kw_short  = kw[:26]
            chip_w    = min(_i(len(kw_short) * 0.075 + 0.22), width - PAD * 2)
            chip_w    = max(chip_w, _i(0.7))
            if chip_x + chip_w > left + width - PAD:
                chip_x  = left + PAD
                chip_y += CHIP_H + _i(0.04)
            _rect(slide, chip_x, chip_y, chip_w, CHIP_H, WHITE(), border=BORDER_CLR())
            _textbox(slide, chip_x + _i(0.04), chip_y + _i(0.03),
                     chip_w - _i(0.06), CHIP_H - _i(0.04),
                     kw_short, size=7, color=MID_GRAY())
            chip_x += chip_w + _i(0.06)
        y = chip_y + CHIP_H + _i(0.08)

    # Description
    if desc and y < top + height - _i(0.22):
        _textbox(slide, left + PAD, y, width - PAD * 2,
                 top + height - _i(0.06) - y,
                 desc, size=7.5, color=MID_GRAY(), wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 4  –  Category-Level Scope  (stacked shapes, like reference deck)
# ─────────────────────────────────────────────────────────────────────────────

def slide4_category_scope(prs, clusters_data: dict, date_str: str,
                           client_name: str, comp_names: list):
    slide = _blank(prs)
    categories: list = clusters_data.get("slide4_categories", [])
    main_comp = comp_names[0] if comp_names else "Competitor"

    _header(slide,
            f"Category-Level Keyword Gap — {client_name} vs {main_comp}",
            f"{client_name}'s ranking gap is structural — deficit exceeds 3x-10x in key categories  |  Source: Ahrefs, {date_str}")
    _footer(slide, f"Source: Ahrefs as of {date_str}")

    if not categories:
        _textbox(slide, _i(0.3), _i(2.5), _i(9.4), _i(0.5),
                 "No category data. Populate slide4_categories in clusters JSON.",
                 size=10, color=MID_GRAY())
        return

    # ── Column definitions (proportions match reference deck) ─────────────────
    ACCENT_W = _i(0.045)
    COL_LEFT = _i(0.22)
    RIGHT_EDGE = _i(9.78)

    col_defs = [
        ("Category",                   _i(1.15)),
        (f"{client_name}\nKWs",        _i(0.52)),
        (f"{client_name}\nTop 10",     _i(0.54)),
        (f"{client_name}\nTraffic",    _i(0.60)),
        (f"{main_comp}\nKWs",          _i(0.52)),
        (f"{main_comp}\nTop 10",       _i(0.54)),
        (f"{main_comp}\nTraffic",      _i(0.60)),
        ("Traffic\nGap",               _i(0.50)),
        ("High-Volume Missing Keywords", None),  # fills remainder
    ]
    fixed_w = sum(w for _, w in col_defs if w is not None)
    remainder = RIGHT_EDGE - COL_LEFT - ACCENT_W - fixed_w
    col_defs[-1] = (col_defs[-1][0], remainder)

    col_lefts = []
    x = COL_LEFT + ACCENT_W
    for _, w in col_defs:
        col_lefts.append(x)
        x += w

    HEADER_TOP = _i(0.62)
    HEADER_H   = _i(0.34)
    ROW_H      = _i(0.50)
    DATA_START = HEADER_TOP + HEADER_H
    TOTAL_W    = RIGHT_EDGE - COL_LEFT

    # ── Column header row ─────────────────────────────────────────────────────
    _rect(slide, COL_LEFT, HEADER_TOP, TOTAL_W, HEADER_H, TEAL())
    for (label, w), lft in zip(col_defs, col_lefts):
        _textbox(slide, lft, HEADER_TOP + _i(0.04), w, HEADER_H - _i(0.06),
                 label, size=7.5, bold=True, color=WHITE(), align="center", wrap=True)

    # ── Data rows ─────────────────────────────────────────────────────────────
    ACCENT_COLORS = [ORANGE(), RED_ALERT(), ORANGE(), RED_ALERT(),
                     ORANGE(), RED_ALERT(), ORANGE(), RED_ALERT()]

    for i, cat in enumerate(categories[:8]):
        row_top = DATA_START + i * ROW_H
        bg = ALT_ROW() if i % 2 == 0 else WHITE()

        _rect(slide, COL_LEFT, row_top, TOTAL_W, ROW_H, bg)
        _rect(slide, COL_LEFT, row_top, ACCENT_W, ROW_H, ACCENT_COLORS[i % 8])

        values = [
            cat.get("category", ""),
            _fmt_num(cat.get("client_kws", 0)),
            _fmt_num(cat.get("client_top10", 0)),
            _fmt_num(cat.get("client_traffic", 0)),
            _fmt_num(cat.get("comp_kws", 0)),
            _fmt_num(cat.get("comp_top10", 0)),
            _fmt_num(cat.get("comp_traffic", 0)),
            None,   # gap badge
            cat.get("missing_keywords", ""),
        ]

        for j, (val, (_, cw), lft) in enumerate(zip(values, col_defs, col_lefts)):
            inner_top = row_top + _i(0.06)
            inner_h   = ROW_H - _i(0.1)
            if j == 0:
                _textbox(slide, lft + _i(0.04), inner_top, cw - _i(0.06), inner_h,
                         str(val), size=8.5, bold=True, color=TEAL())
            elif j == 7:
                gap_str    = cat.get("gap_multiple", "")
                bg_c, fg_c = _gap_colors(gap_str)
                BG_H = _i(0.28)
                BG_W = cw - _i(0.06)
                badge_t = row_top + (ROW_H - BG_H) // 2
                _rect(slide, lft + _i(0.03), badge_t, BG_W, BG_H, bg_c)
                _textbox(slide, lft + _i(0.03), badge_t, BG_W, BG_H,
                         gap_str, size=8.5, bold=True, color=fg_c, align="center")
            elif j == 8:
                _textbox(slide, lft + _i(0.04), inner_top, cw - _i(0.06), inner_h,
                         str(val), size=7, color=MID_GRAY(), wrap=True)
            elif j in (1, 2, 3):
                _textbox(slide, lft, inner_top, cw, inner_h,
                         str(val), size=8.5, bold=(j == 3), color=TEAL(), align="center")
            else:
                _textbox(slide, lft, inner_top, cw, inner_h,
                         str(val), size=8.5, color=MID_GRAY(), align="center")

        _hline(slide, COL_LEFT, row_top + ROW_H, TOTAL_W)


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 5  –  Missing Sub-Category Pages
# ─────────────────────────────────────────────────────────────────────────────

def slide5_onpage_opportunity(prs, clusters_data: dict, date_str: str, client_name: str):
    slide = _blank(prs)
    pages: list = clusters_data.get("slide5_pages", [])

    _header(slide,
            "Missing Sub-Category Pages — Specific On-Page Opportunity Areas",
            f"Competitors earn category traffic through style-specific pages. {client_name} has the products — not the pages.  |  Source: Ahrefs, {date_str}")
    _footer(slide, f"Source: {client_name} website + Ahrefs, {date_str}")

    if not pages:
        _textbox(slide, _i(0.3), _i(2.5), _i(9.4), _i(0.5),
                 "No page data. Populate slide5_pages in clusters JSON.",
                 size=10, color=MID_GRAY())
        return

    n         = min(len(pages), 6)
    COLS      = 3
    ROWS_CNT  = (n + COLS - 1) // COLS
    GAP_X     = _i(0.14)
    GAP_Y     = _i(0.12)
    GRID_LEFT = _i(0.22)
    GRID_TOP  = _i(0.62)
    GRID_W    = _i(9.56)
    GRID_H    = _i(5.44) - GRID_TOP - _i(0.12)
    CARD_W    = (GRID_W - (COLS - 1) * GAP_X) // COLS
    CARD_H    = (GRID_H - (ROWS_CNT - 1) * GAP_Y) // ROWS_CNT

    def _card_colors(status: str):
        sl = status.lower()
        if "0 traffic" in sl or "no " in sl or "does not" in sl:
            return RED_ALERT(), RED_TINT()
        if "wrong" in sl or "retarget" in sl or "wrong slug" in sl:
            return ORANGE(), ORANGE_TINT()
        return TEAL(), TEAL_TINT()

    for idx, pg in enumerate(pages[:n]):
        col_i  = idx % COLS
        row_i  = idx // COLS
        card_l = GRID_LEFT + col_i * (CARD_W + GAP_X)
        card_t = GRID_TOP + row_i * (CARD_H + GAP_Y)

        category = pg.get("category", "")
        status   = pg.get("client_status", "")
        create   = pg.get("pages_to_create", [])
        comp_ex  = pg.get("comp_traffic_examples", [])
        kws      = pg.get("top_keywords", "")

        bar_color, bg_color = _card_colors(status)
        BAR_H = _i(0.18)
        PAD   = _i(0.1)

        _rect(slide, card_l, card_t, CARD_W, CARD_H, bg_color, border=BORDER_CLR())
        _rect(slide, card_l, card_t, CARD_W, BAR_H, bar_color)

        _textbox(slide, card_l + PAD, card_t + BAR_H + _i(0.05),
                 CARD_W - PAD * 2, _i(0.26),
                 category, size=8.5, bold=True, color=bar_color)

        y = card_t + BAR_H + _i(0.34)

        if status:
            _textbox(slide, card_l + PAD, y, CARD_W - PAD * 2, _i(0.28),
                     f"Today: {status}", size=7.5, color=MID_GRAY(), wrap=True)
            y += _i(0.30)

        if create:
            _textbox(slide, card_l + PAD, y, CARD_W - PAD * 2, _i(0.28),
                     "Create:  " + "  /  ".join(create[:3]),
                     size=7.5, color=DARK_NAVY(), wrap=True)
            y += _i(0.30)

        for cex in comp_ex[:2]:
            chip_text = (f"{cex.get('comp', '')}  "
                         f"{cex.get('page', '')}:  "
                         f"{_fmt_num(cex.get('traffic', 0))}/mo")
            CHIP_H = _i(0.20)
            _rect(slide, card_l + PAD, y, CARD_W - PAD * 2, CHIP_H,
                  ORANGE_TINT(), border=BORDER_CLR())
            _textbox(slide, card_l + PAD + _i(0.05), y + _i(0.03),
                     CARD_W - PAD * 2 - _i(0.1), CHIP_H - _i(0.04),
                     chip_text, size=7, color=NEAR_BLACK())
            y += CHIP_H + _i(0.04)

        if kws and y < card_t + CARD_H - _i(0.18):
            _textbox(slide, card_l + PAD, y + _i(0.04),
                     CARD_W - PAD * 2, card_t + CARD_H - _i(0.06) - y,
                     f"KWs: {kws[:85]}", size=7, color=MID_GRAY(), wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
#  Slide 6  –  In-Page Opportunity: Category Deep-Dives
# ─────────────────────────────────────────────────────────────────────────────

def slide6_inpage_opportunity(prs, clusters_data: dict, date_str: str, client_name: str):
    slide = _blank(prs)
    items: list = clusters_data.get("slide6_inpage", [])

    _header(slide,
            "In-Page Opportunity — Category Deep-Dive",
            f"Mis-targeted pages, wrong keyword focus, slug issues — {client_name} state vs competitor benchmark  |  Source: Ahrefs, {date_str}")
    _footer(slide, f"Source: Ahrefs Top Pages Export  |  {date_str}")

    if not items:
        _textbox(slide, _i(0.3), _i(2.5), _i(9.4), _i(0.5),
                 "No in-page data. Populate slide6_inpage in clusters JSON.",
                 size=10, color=MID_GRAY())
        return

    COL_LABEL_W  = _i(1.2)
    COL_CLIENT_W = _i(2.5)
    COL_COMP_W   = _i(2.8)
    TOTAL_W      = _i(9.56)
    COL_ACTION_W = TOTAL_W - COL_LABEL_W - COL_CLIENT_W - COL_COMP_W
    LEFT = _i(0.22)

    col_lefts  = [
        LEFT,
        LEFT + COL_LABEL_W,
        LEFT + COL_LABEL_W + COL_CLIENT_W,
        LEFT + COL_LABEL_W + COL_CLIENT_W + COL_COMP_W,
    ]
    col_widths = [COL_LABEL_W, COL_CLIENT_W, COL_COMP_W, COL_ACTION_W]
    col_hdrs   = ["Category", f"{client_name} Today", "Competitor Benchmark", "Opportunity / Action"]

    HDR_TOP = _i(0.62)
    HDR_H   = _i(0.28)
    _rect(slide, LEFT, HDR_TOP, TOTAL_W, HDR_H, ALT_ROW())
    for lbl, lft, cw in zip(col_hdrs, col_lefts, col_widths):
        _textbox(slide, lft + _i(0.04), HDR_TOP + _i(0.05), cw - _i(0.06), HDR_H - _i(0.08),
                 lbl, size=8, bold=True, color=DARK_NAVY())
    _hline(slide, LEFT, HDR_TOP + HDR_H, TOTAL_W)

    n_items    = min(len(items), 4)
    ROW_GAP    = _i(0.05)
    avail_h    = _i(5.44) - HDR_TOP - HDR_H - ROW_GAP * n_items - _i(0.1)
    ROW_H      = min(avail_h // n_items, _i(1.1))
    DATA_START = HDR_TOP + HDR_H + ROW_GAP

    for i, item in enumerate(items[:n_items]):
        row_top = DATA_START + i * (ROW_H + ROW_GAP)
        PAD     = _i(0.08)
        inner_h = ROW_H - PAD * 2

        # Category — dark navy pill
        _rect(slide, col_lefts[0], row_top, col_widths[0], ROW_H, DARK_NAVY())
        _textbox(slide, col_lefts[0] + PAD, row_top + PAD,
                 col_widths[0] - PAD * 2, inner_h,
                 item.get("category", ""), size=8, bold=True, color=WHITE(), wrap=True)

        # Client today — warm pink tint (signals problem state)
        _rect(slide, col_lefts[1], row_top, col_widths[1], ROW_H,
              WARM_PINK(), border=BORDER_CLR())
        client_lines = [
            f"URL: {item.get('client_url', 'N/A')}",
            f"Traffic: {_fmt_num(item.get('client_traffic', 0))}/mo",
            f"Top KW: \"{item.get('client_top_kw', '')}\" ({_fmt_num(item.get('client_kw_sv', 0))} SV)",
        ]
        issue = item.get("issue", "")
        if issue:
            client_lines.append(f"Issue: {issue}")
        _textbox_lines(slide, col_lefts[1] + PAD, row_top + PAD,
                       col_widths[1] - PAD * 2, inner_h, client_lines, size=7.5)

        # Competitor — teal tint (benchmark / target state)
        _rect(slide, col_lefts[2], row_top, col_widths[2], ROW_H,
              TEAL_TINT(), border=BORDER_CLR())
        comp_lines = [
            f"{item.get('comp_name', 'Comp')} — Traffic: {_fmt_num(item.get('comp_traffic', 0))}/mo",
            f"Top KW: \"{item.get('comp_top_kw', '')}\" ({_fmt_num(item.get('comp_kw_sv', 0))} SV)",
        ]
        _textbox_lines(slide, col_lefts[2] + PAD, row_top + PAD,
                       col_widths[2] - PAD * 2, inner_h,
                       comp_lines, size=7.5, color=TEAL())

        # Action — light gray with a teal "Action" label strip
        _rect(slide, col_lefts[3], row_top, col_widths[3], ROW_H,
              LIGHT_GRAY(), border=BORDER_CLR())
        LABEL_H = _i(0.2)
        _rect(slide, col_lefts[3], row_top, col_widths[3], LABEL_H, TEAL())
        _textbox(slide, col_lefts[3] + PAD, row_top + _i(0.04),
                 col_widths[3] - PAD * 2, LABEL_H - _i(0.05),
                 "Action", size=7, bold=True, color=WHITE())
        _textbox(slide, col_lefts[3] + PAD, row_top + LABEL_H + _i(0.04),
                 col_widths[3] - PAD * 2, inner_h - LABEL_H,
                 item.get("action", ""), size=7.5,
                 color=_rgb(0x1A, 0x73, 0x48), wrap=True)

        _hline(slide, LEFT, row_top + ROW_H, TOTAL_W)


# ─────────────────────────────────────────────────────────────────────────────
#  Entry point
# ─────────────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="seo-proposal-slides generator")
    ap.add_argument("--client-name", required=True)
    ap.add_argument("--client-csv", required=True)
    ap.add_argument("--competitor", action="append", default=[], metavar="NAME:PATH",
                    help="Repeat for each competitor, e.g. 'Senco:/path/senco.csv'")
    ap.add_argument("--top-pages-csv", default=None, help="Client Top Pages CSV from Ahrefs")
    ap.add_argument("--comp-pages", action="append", default=[], metavar="NAME:PATH",
                    help="Competitor top pages CSV, e.g. 'Senco:/path/senco-pages.csv'")
    ap.add_argument("--output-dir", default=None,
                    help="Where to write outputs (default: same dir as --client-csv)")
    ap.add_argument("--clusters-json", default=None,
                    help="Path to Claude's enrichment JSON (required for --finalize)")
    ap.add_argument("--finalize", action="store_true",
                    help="Generate PPTX using --clusters-json")
    args = ap.parse_args()

    client_name = args.client_name
    slug = re.sub(r"[^\w]", "-", client_name.lower())

    if args.output_dir:
        output_dir = Path(args.output_dir).expanduser()
    else:
        output_dir = Path(args.client_csv).expanduser().parent
    output_dir.mkdir(parents=True, exist_ok=True)

    # Parse competitors
    comp_names: list[str] = []
    comp_csvs:  list[str] = []
    for c in args.competitor:
        parts = c.split(":", 1)
        if len(parts) == 2:
            comp_names.append(parts[0])
            comp_csvs.append(parts[1])
        else:
            print(f"WARNING: --competitor '{c}' not in NAME:PATH format — skipping", file=sys.stderr)

    # Parse comp pages
    comp_pages_map: dict[str, str] = {}
    for p in args.comp_pages:
        parts = p.split(":", 1)
        if len(parts) == 2:
            comp_pages_map[parts[0]] = parts[1]

    # ── Load CSVs ────────────────────────────────────────────────────────────
    print(f"Loading {client_name} keywords from {args.client_csv} ...")
    client_df = load_kw(args.client_csv)
    print(f"  → {len(client_df)} rows")

    comp_dfs: list[pd.DataFrame] = []
    for name, path in zip(comp_names, comp_csvs):
        print(f"Loading {name} keywords from {path} ...")
        df = load_kw(path)
        print(f"  → {len(df)} rows")
        comp_dfs.append(df)

    top_pages_map: dict[str, pd.DataFrame] = {}
    if args.top_pages_csv:
        print(f"Loading {client_name} top pages ...")
        top_pages_map[client_name] = load_kw(args.top_pages_csv)
    for name, path in comp_pages_map.items():
        print(f"Loading {name} top pages ...")
        top_pages_map[name] = load_kw(path)

    # ── Pass 1: compute base data ─────────────────────────────────────────────
    print("\nComputing slide 1 (keyword ranking)...")
    s1 = compute_slide1(client_df, comp_dfs, comp_names)
    print(f"  → {len(s1['rows'])} keyword rows")

    print("Computing slide 2 (ranking bucket)...")
    s2 = compute_slide2(client_df, client_name)

    print("Extracting data for Claude (slides 3-6)...")
    for_claude = summarise_for_claude(client_df, comp_dfs, client_name, comp_names, top_pages_map)

    base_json = output_dir / f"proposals-{slug}.json"
    with open(base_json, "w", encoding="utf-8") as f:
        json.dump(
            {
                "client_name": client_name,
                "comp_names": comp_names,
                "slide1": s1,
                "slide2": s2,
                "for_claude": for_claude,
            },
            f,
            ensure_ascii=False,
            indent=2,
        )

    print(f"\n[Pass 1 complete] Written: {base_json}")
    clusters_target = output_dir / f"clusters-{slug}.json"
    if not args.finalize:
        print(f"""
Next steps:
  1. Read: for_claude.client_top_keywords  (top 300 client keywords)
           for_claude.competitor_keywords  (top 200 per competitor)
           for_claude.top_pages            (page-level data)

  2. Produce Claude enrichment JSON at:
       {clusters_target}
     (see SKILL.md for the exact schema — slides 3-6 require it)

  3. Re-run with:
       --clusters-json "{clusters_target}" --finalize
""")
        return

    # ── Pass 2: generate PPTX ─────────────────────────────────────────────────
    if not args.clusters_json:
        print("ERROR: --finalize requires --clusters-json <path>", file=sys.stderr)
        sys.exit(1)

    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    print(f"\nLoading clusters JSON: {args.clusters_json}")
    with open(args.clusters_json, encoding="utf-8") as f:
        clusters = json.load(f)

    date_str = datetime.today().strftime("%B %Y")
    prs = _make_prs()

    print("Generating slides...")
    slide1_keyword_ranking(prs, s1, date_str, client_name)
    print("  [1/6] Keyword ranking")
    slide2_ranking_bucket(prs, s2, date_str)
    print("  [2/6] Ranking buckets")
    slide3_topical_cluster(prs, clusters, date_str, client_name)
    print("  [3/6] Topical clusters")
    slide4_category_scope(prs, clusters, date_str, client_name, comp_names)
    print("  [4/6] Category scope")
    slide5_onpage_opportunity(prs, clusters, date_str, client_name)
    print("  [5/6] On-page opportunities")
    slide6_inpage_opportunity(prs, clusters, date_str, client_name)
    print("  [6/6] In-page deep-dive")

    out_pptx = output_dir / f"proposal-slides-{slug}.pptx"
    prs.save(str(out_pptx))
    print(f"\n[Done] PPTX saved to: {out_pptx}")


if __name__ == "__main__":
    main()
